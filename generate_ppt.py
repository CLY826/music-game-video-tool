#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
音游视频练习辅助小程序 - 期末大作业汇报PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 颜色定义 ====================
THEME = RGBColor(0x6C, 0x63, 0xFF)       # 主色 #6c63ff
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)     # 深色背景 #1a1a2e
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x88)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)  # 绿色强调
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)  # 橙色强调
CARD_BG = RGBColor(0xF8, 0xF8, 0xFC)       # 卡片背景
SUBTLE_PURPLE = RGBColor(0xE8, 0xE6, 0xFF)  # 淡紫色

# ==================== 工具函数 ====================
def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, font_color=DARK_TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=16, font_color=DARK_TEXT, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(4), font_name='Microsoft YaHei'):
    """添加段落"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.5)):
    """添加紫色强调竖条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME
    bar.line.fill.background()
    return bar

def add_bottom_bar(slide):
    """添加底部装饰条"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.1),
        Inches(13.33), Inches(0.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BG
    bar.line.fill.background()
    return bar

def add_page_number(slide, num, total):
    """添加页码"""
    add_textbox(slide, Inches(12.4), Inches(7.15), Inches(0.8), Inches(0.3),
                f"{num}/{total}", font_size=10, font_color=WHITE, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, title, subtitle=""):
    """添加章节标题（左上角）"""
    add_accent_bar(slide, Inches(0.6), Inches(0.4), height=Inches(0.45))
    add_textbox(slide, Inches(0.85), Inches(0.35), Inches(8), Inches(0.6),
                title, font_size=28, font_color=DARK_TEXT, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.85), Inches(0.9), Inches(8), Inches(0.4),
                    subtitle, font_size=14, font_color=MEDIUM_GRAY)

def add_feature_card(slide, left, top, width, height, icon, title, desc):
    """添加功能卡片"""
    card = add_shape(slide, left, top, width, height, fill_color=WHITE)
    card.shadow.inherit = False

    # 图标区域（圆形背景）
    icon_bg = add_shape(slide, left + Inches(0.25), top + Inches(0.25),
                        Inches(0.55), Inches(0.55),
                        fill_color=SUBTLE_PURPLE,
                        shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.28),
                Inches(0.55), Inches(0.5),
                icon, font_size=20, alignment=PP_ALIGN.CENTER)

    # 标题
    add_textbox(slide, left + Inches(0.2), top + Inches(0.9),
                width - Inches(0.4), Inches(0.35),
                title, font_size=14, font_color=DARK_TEXT, bold=True,
                alignment=PP_ALIGN.CENTER)

    # 描述
    add_textbox(slide, left + Inches(0.2), top + Inches(1.2),
                width - Inches(0.4), height - Inches(1.4),
                desc, font_size=10, font_color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)

# ==================== 创建演示文稿 ====================
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

TOTAL_SLIDES = 10

# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_bg(slide1, DARK_BG)

# 装饰圆形
add_shape(slide1, Inches(-1.5), Inches(-1.5), Inches(4), Inches(4),
          fill_color=RGBColor(0x2A, 0x2A, 0x4E), shape_type=MSO_SHAPE.OVAL)
add_shape(slide1, Inches(11), Inches(4.5), Inches(4), Inches(4),
          fill_color=RGBColor(0x2A, 0x2A, 0x4E), shape_type=MSO_SHAPE.OVAL)

# 主标题
add_textbox(slide1, Inches(2), Inches(1.8), Inches(9.33), Inches(1.2),
            "音游视频练习辅助小程序", font_size=44, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

# 英文副标题
add_textbox(slide1, Inches(2), Inches(3.0), Inches(9.33), Inches(0.6),
            "Rhythm Practice — WeChat Mini Program", font_size=20,
            font_color=THEME, alignment=PP_ALIGN.CENTER)

# 分割线
add_shape(slide1, Inches(5.2), Inches(3.8), Inches(2.93), Inches(0.04),
          fill_color=THEME, shape_type=MSO_SHAPE.RECTANGLE)

# 项目信息
add_textbox(slide1, Inches(2), Inches(4.2), Inches(9.33), Inches(0.5),
            "期末大作业汇报", font_size=22, font_color=WHITE,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide1, Inches(2), Inches(5.0), Inches(9.33), Inches(0.4),
            "基于微信小程序 + 云开发 | AI 智能润色 | AB 循环播放 | 变速练习",
            font_size=14, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# 底部信息
add_textbox(slide1, Inches(2), Inches(6.2), Inches(9.33), Inches(0.3),
            "GitHub: github.com/CLY826/music-game-video-tool",
            font_size=12, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide1, 1, TOTAL_SLIDES)

# ============================================================
# 第2页：项目概况
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, LIGHT_GRAY)
add_section_header(slide2, "项目概况", "Project Overview")
add_bottom_bar(slide2)
add_page_number(slide2, 2, TOTAL_SLIDES)

# 左侧：项目描述卡片
desc_card = add_shape(slide2, Inches(0.6), Inches(1.5), Inches(6), Inches(3.5),
                       fill_color=WHITE)
add_textbox(slide2, Inches(0.9), Inches(1.7), Inches(5.4), Inches(0.4),
            "项目简介", font_size=18, font_color=THEME, bold=True)

desc_box = add_textbox(slide2, Inches(0.9), Inches(2.2), Inches(5.4), Inches(2.5),
                        "", font_size=13, font_color=DARK_TEXT)
tf = desc_box.text_frame
tf.paragraphs[0].text = "「音游视频练习辅助小程序」是一款专为音游玩家设计的练习辅助工具。"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.name = 'Microsoft YaHei'
tf.paragraphs[0].font.color.rgb = DARK_TEXT

descriptions = [
    "玩家可以上传自己的练习视频，按音游分类浏览和管理，",
    "利用 AB 循环和变速播放反复练习难点段落，",
    "还能借助 AI 润色功能优化备注说明，让分享更专业。",
]
for d in descriptions:
    add_paragraph(tf, d, font_size=13)

add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "技术栈：微信原生小程序 + 云开发 + TokenHub AI", font_size=12, font_color=THEME, bold=True)

# 右侧：关键数据
metrics_card = add_shape(slide2, Inches(7), Inches(1.5), Inches(5.7), Inches(3.5),
                          fill_color=DARK_BG)

# 数据项
metrics = [
    ("7", "功能页面", "首页/分类/上传/播放/搜索/社区/个人中心"),
    ("6", "云函数", "login / getVideos / uploadVideo / addComment / getComments / aiPolish"),
    ("4", "数据库集合", "game / video / comment / user"),
    ("97%+", "后端覆盖率", "Jest 单元测试 + 接口测试"),
]

for i, (num, label, detail) in enumerate(metrics):
    y_pos = Inches(1.75) + Inches(i * 0.8)
    add_textbox(slide2, Inches(7.4), y_pos, Inches(1.2), Inches(0.5),
                num, font_size=28, font_color=THEME, bold=True)
    add_textbox(slide2, Inches(8.6), y_pos, Inches(2), Inches(0.3),
                label, font_size=14, font_color=WHITE, bold=True)
    add_textbox(slide2, Inches(8.6), y_pos + Inches(0.28), Inches(3.8), Inches(0.3),
                detail, font_size=10, font_color=MEDIUM_GRAY)

# 底部功能标签
features_row = ["游戏分类", "视频上传", "AB 循环播放", "变速播放", "关键词搜索", "社区互动", "AI 润色"]
for i, feat in enumerate(features_row):
    x = Inches(0.6) + Inches(i * 1.75)
    tag = add_shape(slide2, x, Inches(5.5), Inches(1.55), Inches(0.4),
                    fill_color=SUBTLE_PURPLE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide2, x, Inches(5.52), Inches(1.55), Inches(0.35),
                feat, font_size=11, font_color=THEME, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第3页：技术架构
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, LIGHT_GRAY)
add_section_header(slide3, "技术架构", "Technical Architecture")
add_bottom_bar(slide3)
add_page_number(slide3, 3, TOTAL_SLIDES)

# 前端层
fe_card = add_shape(slide3, Inches(0.6), Inches(1.5), Inches(12.13), Inches(1.4),
                     fill_color=WHITE)
add_textbox(slide3, Inches(0.9), Inches(1.6), Inches(2), Inches(0.35),
            "前端展示层", font_size=16, font_color=THEME, bold=True)

fe_items = [
    ("微信小程序原生框架", "WXML / WXSS / JS"),
    ("7 个功能页面", "首页 | 分类 | 上传 | 播放 | 搜索 | 社区 | 个人中心"),
    ("Tab Bar 导航", "首页 / 社区 / 我的"),
]
for i, (title, sub) in enumerate(fe_items):
    x = Inches(0.9) + Inches(i * 3.9)
    add_textbox(slide3, x, Inches(2.05), Inches(3.6), Inches(0.3),
                title, font_size=12, font_color=DARK_TEXT, bold=True)
    add_textbox(slide3, x, Inches(2.35), Inches(3.6), Inches(0.3),
                sub, font_size=10, font_color=MEDIUM_GRAY)

# 箭头
add_textbox(slide3, Inches(5.5), Inches(2.9), Inches(2.33), Inches(0.4),
            "⬇  云开发 SDK  ⬇", font_size=14, font_color=THEME,
            bold=True, alignment=PP_ALIGN.CENTER)

# 云开发层
cloud_card = add_shape(slide3, Inches(0.6), Inches(3.3), Inches(12.13), Inches(1.8),
                        fill_color=WHITE)
add_textbox(slide3, Inches(0.9), Inches(3.4), Inches(3), Inches(0.35),
            "微信云开发层", font_size=16, font_color=THEME, bold=True)

# 三个子模块
cloud_modules = [
    ("云函数", "login\ngetVideos\nuploadVideo\naddComment\ngetComments\naiPolish"),
    ("云数据库", "game 集合\nvideo 集合\ncomment 集合\nuser 集合"),
    ("云存储", "视频文件存储\n封面图存储\n自动上传\n临时 URL"),
]
for i, (title, content) in enumerate(cloud_modules):
    x = Inches(0.9) + Inches(i * 4)
    # 小卡片
    sub_card = add_shape(slide3, x, Inches(3.85), Inches(3.5), Inches(1.1),
                          fill_color=CARD_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide3, x + Inches(0.15), Inches(3.9), Inches(1.2), Inches(0.3),
                title, font_size=12, font_color=THEME, bold=True)
    add_textbox(slide3, x + Inches(1.3), Inches(3.88), Inches(2), Inches(1),
                content, font_size=9, font_color=DARK_TEXT)

# 箭头
add_textbox(slide3, Inches(5.5), Inches(5.1), Inches(2.33), Inches(0.4),
            "⬇  HTTPS API  ⬇", font_size=14, font_color=THEME,
            bold=True, alignment=PP_ALIGN.CENTER)

# AI 服务层
ai_card = add_shape(slide3, Inches(0.6), Inches(5.5), Inches(12.13), Inches(1.2),
                     fill_color=WHITE)
add_textbox(slide3, Inches(0.9), Inches(5.6), Inches(3), Inches(0.35),
            "AI 服务层", font_size=16, font_color=THEME, bold=True)

ai_details = [
    ("TokenHub API", "大模型 API 平台"),
    ("Hy3-preview", "文本生成模型"),
    ("OpenAI 兼容格式", "标准 Chat Completions 接口"),
    ("智能润色", "音游社区写作助手"),
]
for i, (title, sub) in enumerate(ai_details):
    x = Inches(0.9) + Inches(i * 3)
    add_textbox(slide3, x, Inches(6.0), Inches(2.5), Inches(0.3),
                title, font_size=12, font_color=DARK_TEXT, bold=True)
    add_textbox(slide3, x, Inches(6.28), Inches(2.5), Inches(0.3),
                sub, font_size=10, font_color=MEDIUM_GRAY)

# ============================================================
# 第4页：核心功能 - 播放器
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, LIGHT_GRAY)
add_section_header(slide4, "核心功能：视频播放器", "Core Feature: Video Player")
add_bottom_bar(slide4)
add_page_number(slide4, 4, TOTAL_SLIDES)

# 左侧 - 播放器功能
player_card = add_shape(slide4, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2),
                         fill_color=WHITE)
add_textbox(slide4, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.4),
            "播放器核心能力", font_size=20, font_color=THEME, bold=True)

player_features = [
    ("🎵  AB 循环播放", "设定起始点(A)和结束点(B)，片段自动循环，\n反复练习难点段落，无需手动拖拽进度条"),
    ("⚡  变速播放", "支持 0.5x / 0.75x / 1.0x / 1.25x / 1.5x 倍速，\n慢放观察细节，正常速度验证手感"),
    ("📋  视频信息展示", "显示歌曲名、音游分类、练习备注、时长等信息，\n配合 AI 润色后的专业描述"),
    ("💬  评论区", "视频下方评论互动，分享练习心得和技巧，\n支持实时获取和发布评论"),
]

for i, (title, desc) in enumerate(player_features):
    y = Inches(2.3) + Inches(i * 1.15)
    add_accent_bar(slide4, Inches(0.9), y, height=Inches(0.7))
    add_textbox(slide4, Inches(1.15), y - Inches(0.05), Inches(5), Inches(0.3),
                title, font_size=14, font_color=DARK_TEXT, bold=True)
    add_textbox(slide4, Inches(1.15), y + Inches(0.28), Inches(5), Inches(0.6),
                desc, font_size=11, font_color=MEDIUM_GRAY)

# 右侧 - 操作流程
flow_card = add_shape(slide4, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2),
                       fill_color=DARK_BG)
add_textbox(slide4, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4),
            "播放器交互流程", font_size=20, font_color=WHITE, bold=True)

flow_steps = [
    ("1", "进入播放页", "通过视频列表/社区/搜索进入"),
    ("2", "加载视频", "自动获取临时 URL 并播放"),
    ("3", "设定 AB 点", "点击标记按钮设定循环范围"),
    ("4", "调整速度", "切换倍速适配练习需求"),
    ("5", "循环播放", "AB 区间自动循环，专注练习"),
    ("6", "查看/发布评论", "互动交流，分享心得"),
]

for i, (num, step, detail) in enumerate(flow_steps):
    y = Inches(2.35) + Inches(i * 0.72)
    # 序号圆圈
    circle = add_shape(slide4, Inches(7.1), y, Inches(0.4), Inches(0.4),
                        fill_color=THEME, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide4, Inches(7.1), y + Inches(0.02), Inches(0.4), Inches(0.35),
                num, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 连接线
    if i < len(flow_steps) - 1:
        add_shape(slide4, Inches(7.28), y + Inches(0.4), Inches(0.04), Inches(0.32),
                  fill_color=THEME, shape_type=MSO_SHAPE.RECTANGLE)
    # 步骤信息
    add_textbox(slide4, Inches(7.7), y, Inches(2), Inches(0.3),
                step, font_size=13, font_color=WHITE, bold=True)
    add_textbox(slide4, Inches(7.7), y + Inches(0.27), Inches(4.5), Inches(0.3),
                detail, font_size=10, font_color=MEDIUM_GRAY)

# ============================================================
# 第5页：AI 润色功能
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, LIGHT_GRAY)
add_section_header(slide5, "AI 润色功能", "AI Polish Feature")
add_bottom_bar(slide5)
add_page_number(slide5, 5, TOTAL_SLIDES)

# 左侧 - AI 工作流程
ai_flow_card = add_shape(slide5, Inches(0.6), Inches(1.5), Inches(6), Inches(5.2),
                          fill_color=WHITE)
add_textbox(slide5, Inches(0.9), Inches(1.7), Inches(5.4), Inches(0.4),
            "AI 润色工作流程", font_size=20, font_color=THEME, bold=True)

# 流程卡片
ai_steps = [
    ("用户输入", "输入练习备注文本（歌曲名+音游名）"),
    ("云函数调用", "aiPolish 云函数接收请求，拼接 Prompt"),
    ("TokenHub API", "发送至 Hy3-preview 模型，temperature=0.7"),
    ("智能润色", "保持原意、优化表达、加入音游术语"),
    ("返回结果", "返回润色后文本，用户确认后保存"),
]

for i, (step, desc) in enumerate(ai_steps):
    y = Inches(2.3) + Inches(i * 0.85)
    # 步骤卡片
    step_card = add_shape(slide5, Inches(0.9), y, Inches(5.2), Inches(0.7),
                           fill_color=CARD_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    # 序号
    num_shape = add_shape(slide5, Inches(1.05), y + Inches(0.1), Inches(0.45), Inches(0.45),
                           fill_color=THEME, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide5, Inches(1.05), y + Inches(0.12), Inches(0.45), Inches(0.4),
                str(i+1), font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, Inches(1.7), y + Inches(0.08), Inches(2), Inches(0.3),
                step, font_size=13, font_color=DARK_TEXT, bold=True)
    add_textbox(slide5, Inches(1.7), y + Inches(0.35), Inches(4.2), Inches(0.3),
                desc, font_size=10, font_color=MEDIUM_GRAY)

# 右侧 - Prompt 设计与示例
prompt_card = add_shape(slide5, Inches(7), Inches(1.5), Inches(5.7), Inches(2.4),
                         fill_color=WHITE)
add_textbox(slide5, Inches(7.3), Inches(1.6), Inches(5.1), Inches(0.4),
            "System Prompt 设计要点", font_size=16, font_color=THEME, bold=True)

prompt_rules = [
    "保持用户原意不变，只优化表达",
    "语气自然、像真实玩家，不官方",
    "加入音游术语（FC、AP、PM、收歌、手元等）",
    "保留关键信息（难点位置、练习目标）",
    "字数控制在 50~150 字",
    "适当使用 emoji 增加趣味性",
]
rules_box = add_textbox(slide5, Inches(7.3), Inches(2.1), Inches(5.1), Inches(1.7),
                         "", font_size=11)
tf = rules_box.text_frame
for j, rule in enumerate(prompt_rules):
    if j == 0:
        tf.paragraphs[0].text = f"  {j+1}. {rule}"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.name = 'Microsoft YaHei'
        tf.paragraphs[0].font.color.rgb = DARK_TEXT
    else:
        add_paragraph(tf, f"  {j+1}. {rule}", font_size=11, space_before=Pt(2), space_after=Pt(2))

# 示例卡片
example_card = add_shape(slide5, Inches(7), Inches(4.15), Inches(5.7), Inches(2.55),
                          fill_color=WHITE)
add_textbox(slide5, Inches(7.3), Inches(4.25), Inches(5.1), Inches(0.4),
            "润色效果示例", font_size=16, font_color=THEME, bold=True)

# 原文
add_shape(slide5, Inches(7.3), Inches(4.75), Inches(5.1), Inches(0.85),
          fill_color=CARD_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide5, Inches(7.5), Inches(4.8), Inches(1), Inches(0.25),
            "原文：", font_size=10, font_color=ACCENT_ORANGE, bold=True)
add_textbox(slide5, Inches(7.5), Inches(5.05), Inches(4.7), Inches(0.5),
            "这个谱面中间那段纵连很难，总是断，多练练",
            font_size=11, font_color=DARK_TEXT)

# 润色后
add_shape(slide5, Inches(7.3), Inches(5.75), Inches(5.1), Inches(0.85),
          fill_color=SUBTLE_PURPLE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide5, Inches(7.5), Inches(5.8), Inches(1), Inches(0.25),
            "润色：", font_size=10, font_color=ACCENT_GREEN, bold=True)
add_textbox(slide5, Inches(7.5), Inches(6.05), Inches(4.7), Inches(0.5),
            "🎵 中段纵连段是 PM 的关键障碍，需要反复练习收歌节奏感，目标先稳定 FC 再冲 AP！💪",
            font_size=11, font_color=DARK_TEXT)

# ============================================================
# 第6页：数据库设计
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, LIGHT_GRAY)
add_section_header(slide6, "数据库设计", "Database Design")
add_bottom_bar(slide6)
add_page_number(slide6, 6, TOTAL_SLIDES)

# 4个集合卡片
collections = [
    {
        "name": "game", "desc": "音游分类",
        "fields": ["name - 音游名称", "icon - Emoji图标", "color - 主题色", "sort - 排序权重", "videoCount - 视频数"],
        "permission": "所有人可读，认证用户可写"
    },
    {
        "name": "video", "desc": "视频信息",
        "fields": ["songName - 歌曲名", "gameId/gameName - 关联音游", "desc - 练习备注", "videoFileId - 云存储ID", "openid - 上传者"],
        "permission": "所有人可读，创建者可写"
    },
    {
        "name": "comment", "desc": "评论数据",
        "fields": ["videoId - 关联视频", "content - 评论内容", "openid - 评论者", "authorName - 昵称", "createTime - 时间"],
        "permission": "所有人可读，创建者可写"
    },
    {
        "name": "user", "desc": "用户信息",
        "fields": ["openid - 唯一标识", "nickName - 昵称", "avatarUrl - 头像", "createTime - 注册时间", "lastLoginTime - 登录时间"],
        "permission": "仅创建者可读写"
    },
]

for i, col in enumerate(collections):
    x = Inches(0.6) + Inches(i * 3.15)
    # 卡片
    card = add_shape(slide6, x, Inches(1.5), Inches(2.9), Inches(4.5),
                      fill_color=WHITE)
    # 标题区
    add_shape(slide6, x, Inches(1.5), Inches(2.9), Inches(0.9),
              fill_color=DARK_BG)
    add_textbox(slide6, x + Inches(0.2), Inches(1.6), Inches(2.5), Inches(0.35),
                col["name"], font_size=20, font_color=THEME, bold=True, font_name='Consolas')
    add_textbox(slide6, x + Inches(0.2), Inches(1.95), Inches(2.5), Inches(0.3),
                col["desc"], font_size=12, font_color=MEDIUM_GRAY)
    # 字段列表
    for j, field in enumerate(col["fields"]):
        y = Inches(2.65) + Inches(j * 0.45)
        add_textbox(slide6, x + Inches(0.2), y, Inches(2.5), Inches(0.35),
                    f"• {field}", font_size=10, font_color=DARK_TEXT)
    # 权限
    add_shape(slide6, x + Inches(0.15), Inches(5.1), Inches(2.6), Inches(0.5),
              fill_color=SUBTLE_PURPLE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide6, x + Inches(0.2), Inches(5.15), Inches(2.5), Inches(0.4),
                col["permission"], font_size=9, font_color=THEME, alignment=PP_ALIGN.CENTER)

# 底部提示
add_textbox(slide6, Inches(0.6), Inches(6.3), Inches(12), Inches(0.3),
            "安全规则示例：{\"read\": true, \"write\": \"auth.openid == doc.openid\"}  — 基于用户身份的细粒度权限控制",
            font_size=10, font_color=MEDIUM_GRAY)

# ============================================================
# 第7页：测试体系
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, LIGHT_GRAY)
add_section_header(slide7, "测试体系", "Testing & Quality")
add_bottom_bar(slide7)
add_page_number(slide7, 7, TOTAL_SLIDES)

# 后端测试卡片
be_card = add_shape(slide7, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.0),
                     fill_color=WHITE)
add_textbox(slide7, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.4),
            "后端测试", font_size=20, font_color=THEME, bold=True)

# 覆盖率大数字
add_textbox(slide7, Inches(0.9), Inches(2.3), Inches(2.5), Inches(0.8),
            "97.01%", font_size=42, font_color=ACCENT_GREEN, bold=True)
add_textbox(slide7, Inches(0.9), Inches(3.1), Inches(2.5), Inches(0.3),
            "行覆盖率", font_size=12, font_color=MEDIUM_GRAY)

# 后端测试详情
be_details = [
    ("单元测试", "30 个", "云函数核心逻辑测试"),
    ("接口测试", "13 个", "API 入参/出参/异常测试"),
    ("Mock 策略", "wx-server-sdk", "模拟云开发 SDK 调用"),
    ("覆盖率工具", "Jest + v8 provider", "text / lcov / html 报告"),
]
for i, (label, value, desc) in enumerate(be_details):
    y = Inches(3.6) + Inches(i * 0.45)
    add_textbox(slide7, Inches(0.9), y, Inches(1.5), Inches(0.3),
                label, font_size=11, font_color=DARK_TEXT, bold=True)
    add_textbox(slide7, Inches(2.4), y, Inches(1.2), Inches(0.3),
                value, font_size=11, font_color=THEME, bold=True)
    add_textbox(slide7, Inches(3.6), y, Inches(2.5), Inches(0.3),
                desc, font_size=10, font_color=MEDIUM_GRAY)

# 前端测试卡片
fe_test_card = add_shape(slide7, Inches(6.8), Inches(1.5), Inches(5.9), Inches(4.0),
                          fill_color=WHITE)
add_textbox(slide7, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4),
            "前端测试", font_size=20, font_color=THEME, bold=True)

# 覆盖率大数字
add_textbox(slide7, Inches(7.1), Inches(2.3), Inches(2.5), Inches(0.8),
            "92.89%", font_size=42, font_color=ACCENT_GREEN, bold=True)
add_textbox(slide7, Inches(7.1), Inches(3.1), Inches(2.5), Inches(0.3),
            "行覆盖率", font_size=12, font_color=MEDIUM_GRAY)

# 前端测试详情
fe_details = [
    ("组件测试", "17 个", "页面生命周期与交互逻辑"),
    ("网络请求测试", "6 个", "云函数调用模拟测试"),
    ("Mock 策略", "wx-globals", "全局 wx 对象模拟"),
    ("覆盖率工具", "Jest + v8 provider", "text / lcov / html 报告"),
]
for i, (label, value, desc) in enumerate(fe_details):
    y = Inches(3.6) + Inches(i * 0.45)
    add_textbox(slide7, Inches(7.1), y, Inches(1.5), Inches(0.3),
                label, font_size=11, font_color=DARK_TEXT, bold=True)
    add_textbox(slide7, Inches(8.6), y, Inches(1.2), Inches(0.3),
                value, font_size=11, font_color=THEME, bold=True)
    add_textbox(slide7, Inches(9.8), y, Inches(2.5), Inches(0.3),
                desc, font_size=10, font_color=MEDIUM_GRAY)

# 覆盖率进度条
add_textbox(slide7, Inches(0.6), Inches(5.8), Inches(5.8), Inches(0.3),
            "后端覆盖率", font_size=12, font_color=DARK_TEXT, bold=True)
bar_bg = add_shape(slide7, Inches(0.6), Inches(6.15), Inches(12.13), Inches(0.25),
                    fill_color=RGBColor(0xE0, 0xE0, 0xE0), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
bar_fill = add_shape(slide7, Inches(0.6), Inches(6.15), Inches(12.13 * 0.97), Inches(0.25),
                      fill_color=ACCENT_GREEN, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

add_textbox(slide7, Inches(0.6), Inches(6.45), Inches(12.13), Inches(0.3),
            "前端覆盖率", font_size=12, font_color=DARK_TEXT, bold=True)
bar_bg2 = add_shape(slide7, Inches(0.6), Inches(6.8), Inches(12.13), Inches(0.25),
                     fill_color=RGBColor(0xE0, 0xE0, 0xE0), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
bar_fill2 = add_shape(slide7, Inches(0.6), Inches(6.8), Inches(12.13 * 0.9289), Inches(0.25),
                       fill_color=THEME, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

# ============================================================
# 第8页：CI/CD 流水线
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, LIGHT_GRAY)
add_section_header(slide8, "CI/CD 流水线", "Continuous Integration & Delivery")
add_bottom_bar(slide8)
add_page_number(slide8, 8, TOTAL_SLIDES)

# 后端 Job 流程
add_textbox(slide8, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4),
            "Backend Job", font_size=18, font_color=THEME, bold=True)

be_steps = [
    ("actions/checkout@v4", "拉取代码"),
    ("actions/setup-node@v4", "Node.js 22 LTS"),
    ("npm ci", "安装根依赖"),
    ("npm ci (tests/)", "安装测试依赖"),
    ("ESLint", "云函数代码规范检查"),
    ("Jest", "运行后端测试 + 覆盖率"),
    ("Codecov", "上传覆盖率报告"),
]

for i, (step, desc) in enumerate(be_steps):
    y = Inches(1.9) + Inches(i * 0.6)
    # 步骤卡片
    step_shape = add_shape(slide8, Inches(0.8), y, Inches(5.2), Inches(0.48),
                            fill_color=WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide8, Inches(1.0), y + Inches(0.06), Inches(2.5), Inches(0.3),
                step, font_size=11, font_color=DARK_TEXT, bold=True, font_name='Consolas')
    add_textbox(slide8, Inches(3.5), y + Inches(0.06), Inches(2.3), Inches(0.3),
                desc, font_size=10, font_color=MEDIUM_GRAY)
    # 连接线
    if i < len(be_steps) - 1:
        add_shape(slide8, Inches(3.4), y + Inches(0.48), Inches(0.03), Inches(0.12),
                  fill_color=THEME, shape_type=MSO_SHAPE.RECTANGLE)

# 前端 Job 流程
add_textbox(slide8, Inches(6.8), Inches(1.4), Inches(6), Inches(0.4),
            "Frontend Job", font_size=18, font_color=THEME, bold=True)

fe_steps = [
    ("actions/checkout@v4", "拉取代码"),
    ("actions/setup-node@v4", "Node.js 22 LTS"),
    ("npm ci", "安装根依赖"),
    ("npm ci (tests/)", "安装测试依赖"),
    ("ESLint", "页面代码规范检查"),
    ("Jest", "运行前端测试 + 覆盖率"),
    ("Codecov", "上传覆盖率报告"),
]

for i, (step, desc) in enumerate(fe_steps):
    y = Inches(1.9) + Inches(i * 0.6)
    step_shape = add_shape(slide8, Inches(7.0), y, Inches(5.2), Inches(0.48),
                            fill_color=WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide8, Inches(7.2), y + Inches(0.06), Inches(2.5), Inches(0.3),
                step, font_size=11, font_color=DARK_TEXT, bold=True, font_name='Consolas')
    add_textbox(slide8, Inches(9.7), y + Inches(0.06), Inches(2.3), Inches(0.3),
                desc, font_size=10, font_color=MEDIUM_GRAY)
    if i < len(fe_steps) - 1:
        add_shape(slide8, Inches(9.6), y + Inches(0.48), Inches(0.03), Inches(0.12),
                  fill_color=THEME, shape_type=MSO_SHAPE.RECTANGLE)

# 触发条件
trigger_card = add_shape(slide8, Inches(0.6), Inches(6.2), Inches(12.13), Inches(0.6),
                          fill_color=DARK_BG)
add_textbox(slide8, Inches(0.9), Inches(6.3), Inches(3), Inches(0.3),
            "触发条件：", font_size=12, font_color=THEME, bold=True)
add_textbox(slide8, Inches(2.5), Inches(6.3), Inches(9), Inches(0.3),
            "Push → main / develop 分支  |  Pull Request → main 分支  |  环境变量: FORCE_JAVASCRIPT_ACTIONS_TO_NODE24=true",
            font_size=11, font_color=WHITE)

# ============================================================
# 第9页：项目结构
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, LIGHT_GRAY)
add_section_header(slide9, "项目结构", "Project Structure")
add_bottom_bar(slide9)
add_page_number(slide9, 9, TOTAL_SLIDES)

# 项目结构卡片
struct_card = add_shape(slide9, Inches(0.6), Inches(1.5), Inches(5.5), Inches(5.2),
                         fill_color=DARK_BG)
add_textbox(slide9, Inches(0.9), Inches(1.65), Inches(5), Inches(0.4),
            "目录结构", font_size=18, font_color=WHITE, bold=True)

tree_text = """rhythm-practice/
├── cloudfunctions/
│   ├── addComment/
│   ├── aiPolish/
│   ├── getComments/
│   ├── getVideos/
│   ├── login/
│   └── uploadVideo/
├── pages/
│   ├── index/       首页
│   ├── category/    分类
│   ├── upload/      上传
│   ├── player/      播放器
│   ├── search/      搜索
│   ├── community/   社区
│   └── profile/     个人中心
├── tests/
│   ├── backend/     后端测试
│   └── frontend/    前端测试
├── .github/
│   └── workflows/
│       └── ci.yml   CI配置
├── app.js / app.json / app.wxss
└── project.config.json"""

add_textbox(slide9, Inches(0.9), Inches(2.1), Inches(5), Inches(4.5),
            tree_text, font_size=10, font_color=RGBColor(0xA0, 0xFF, 0xA0),
            font_name='Consolas')

# 右侧：技术选型
tech_card = add_shape(slide9, Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.2),
                       fill_color=WHITE)
add_textbox(slide9, Inches(6.8), Inches(1.65), Inches(5.6), Inches(0.4),
            "技术选型总览", font_size=18, font_color=THEME, bold=True)

tech_items = [
    ("前端框架", "微信小程序原生开发", "WXML / WXSS / JavaScript，无需额外框架"),
    ("后端服务", "微信云开发", "云函数 + 云数据库 + 云存储，零运维"),
    ("AI 服务", "TokenHub API (Hy3-preview)", "OpenAI 兼容接口，智能文本润色"),
    ("测试框架", "Jest + ESLint", "单元测试、组件测试、代码规范检查"),
    ("CI/CD", "GitHub Actions + Codecov", "自动化构建、测试、覆盖率追踪"),
    ("运行时", "Node.js 22 LTS", "CI 环境与云函数运行时"),
    ("版本控制", "Git + GitHub", "分支管理、PR 审查、持续集成"),
]

for i, (label, value, desc) in enumerate(tech_items):
    y = Inches(2.2) + Inches(i * 0.65)
    add_accent_bar(slide9, Inches(6.8), y + Inches(0.05), height=Inches(0.4))
    add_textbox(slide9, Inches(7.05), y, Inches(1.5), Inches(0.3),
                label, font_size=12, font_color=DARK_TEXT, bold=True)
    add_textbox(slide9, Inches(8.5), y, Inches(3.8), Inches(0.3),
                value, font_size=12, font_color=THEME, bold=True)
    add_textbox(slide9, Inches(8.5), y + Inches(0.25), Inches(3.8), Inches(0.3),
                desc, font_size=9, font_color=MEDIUM_GRAY)

# ============================================================
# 第10页：总结与展望
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, DARK_BG)

# 装饰
add_shape(slide10, Inches(10.5), Inches(-1), Inches(4), Inches(4),
          fill_color=RGBColor(0x2A, 0x2A, 0x4E), shape_type=MSO_SHAPE.OVAL)

add_textbox(slide10, Inches(2), Inches(0.8), Inches(9.33), Inches(0.8),
            "总结与展望", font_size=36, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

# 分割线
add_shape(slide10, Inches(5.2), Inches(1.6), Inches(2.93), Inches(0.04),
          fill_color=THEME, shape_type=MSO_SHAPE.RECTANGLE)

# 项目成果
add_textbox(slide10, Inches(1.5), Inches(2.0), Inches(10.33), Inches(0.4),
            "项目成果", font_size=22, font_color=THEME, bold=True, alignment=PP_ALIGN.CENTER)

achievements = [
    "完成 7 个功能页面的设计与开发，覆盖完整的用户使用流程",
    "实现 6 个云函数，支持登录、视频管理、评论、AI 润色等核心功能",
    "后端测试覆盖率 97.01%，前端测试覆盖率 92.89%，总计 66 个测试用例",
    "搭建 GitHub Actions CI/CD 流水线，自动化代码质量检查与覆盖率追踪",
    "集成 TokenHub AI API，实现智能备注润色，提升用户体验",
]
for i, ach in enumerate(achievements):
    y = Inches(2.5) + Inches(i * 0.45)
    add_textbox(slide10, Inches(1.5), y, Inches(10.33), Inches(0.4),
                f"✓  {ach}", font_size=13, font_color=WHITE)

# 未来展望
add_textbox(slide10, Inches(1.5), Inches(4.9), Inches(10.33), Inches(0.4),
            "未来展望", font_size=22, font_color=THEME, bold=True, alignment=PP_ALIGN.CENTER)

future_items = [
    "Docker 容器化部署，提升环境一致性与部署效率",
    "安全审查与加固，完善数据权限和 API 安全策略",
    "云服务正式部署，支持生产环境运行",
    "监控与日志系统，实时追踪服务状态与异常",
    "更多 AI 功能：自动标签、谱面分析、练习建议",
]
for i, item in enumerate(future_items):
    y = Inches(5.4) + Inches(i * 0.38)
    add_textbox(slide10, Inches(1.5), y, Inches(10.33), Inches(0.35),
                f"→  {item}", font_size=12, font_color=MEDIUM_GRAY)

# 底部感谢
add_textbox(slide10, Inches(2), Inches(7.0), Inches(9.33), Inches(0.4),
            "感谢聆听！  |  GitHub: github.com/CLY826/music-game-video-tool",
            font_size=12, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ==================== 保存文件 ====================
output_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "音游视频练习辅助小程序_期末汇报.pptx"
)
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
